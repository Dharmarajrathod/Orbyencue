import io
import os
import re
import shutil
import subprocess
import threading
import time
import tempfile
import json
import logging
import wave
from collections import deque
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

try:
    from config import load_environment
    from config import ensure_vosk_model_path
except ImportError:
    load_environment = None
    ensure_vosk_model_path = None

from processing_pipeline import clean_transcript, is_meaningful_question_or_request


if load_environment:
    load_environment()


class AnswerRequest(BaseModel):
    question: str = Field(..., min_length=1, max_length=20000)


class AnswerResponse(BaseModel):
    answer: str
    model: str


class KnowledgeChunk(BaseModel):
    text: str
    filename: str
    page: Optional[int] = None


class KnowledgeResponse(BaseModel):
    chunks: list[KnowledgeChunk]
    chunkCount: int
    filename: str


class TranscriptResponse(BaseModel):
    transcript: str
    model: str
    isFinal: bool = True
    chunkNumber: int = 0
    chunkSize: int = 0
    duration: float = 0.0
    sampleRate: int = 0
    audioEnergy: float = 0.0
    voiceActivity: bool = False
    language: str = "unknown"
    languageConfidence: float = 0.0
    transcriptionConfidence: float = 0.0
    meaningful: bool = False
    discarded: bool = False
    reason: str = ""
    latencyMs: int = 0
    stageTimings: dict[str, float] = Field(default_factory=dict)


app = FastAPI(title="ORBYNECUE AI Backend")
BASE_DIR = Path(__file__).resolve().parent
SUPPORTED_UPLOAD_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".pdf", ".doc", ".docx", ".ppt", ".pptx"}
AUDIO_DEBUG_LOGS = deque(maxlen=300)
logger = logging.getLogger("orbynecue.audio")
logging.basicConfig(level=logging.INFO)
STREAMING_STT_SESSIONS = {}
STREAMING_STT_LOCK = threading.Lock()
STREAMING_STT_TTL_SECONDS = 15 * 60
STREAMING_VOSK_MODEL = None
GEMINI_STT_CLIENT = None
DEFAULT_GEMINI_STT_MODEL = "gemini-2.5-flash-lite"


def get_allowed_origins():
    configured = os.getenv("ALLOWED_ORIGINS")
    if configured:
        return [origin.strip() for origin in configured.split(",") if origin.strip()]
    return ["*"]


app.add_middleware(
    CORSMiddleware,
    allow_origins=get_allowed_origins(),
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)

app.mount("/assets", StaticFiles(directory=BASE_DIR / "assets"), name="assets")


def stream_chunks(text: str, max_words: int = 220):
    paragraphs = [paragraph.strip() for paragraph in text.splitlines() if paragraph.strip()]
    buffer = []
    word_count = 0

    for paragraph in paragraphs:
        words = paragraph.split()
        if not words:
            continue
        if word_count + len(words) > max_words and buffer:
            yield "\n".join(buffer)
            buffer = []
            word_count = 0
        buffer.append(paragraph)
        word_count += len(words)

    if buffer:
        yield "\n".join(buffer)


def normalize_text(text: str) -> str:
    text = re.sub(r"\n{2,}", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def extract_upload_text(path: Path, extension: str) -> str:
    return "\n".join(section["text"] for section in extract_upload_sections(path, extension))


def extract_upload_sections(path: Path, extension: str) -> list[dict]:
    if extension in {".txt", ".md", ".json"}:
        return [{"text": path.read_text(encoding="utf-8", errors="ignore"), "page": None}]

    if extension == ".csv":
        import csv

        rows = []
        with path.open(newline="", encoding="utf-8", errors="ignore") as file_obj:
            for row in csv.reader(file_obj):
                rows.append(" ".join(cell for cell in row if cell))
        return [{"text": "\n".join(rows), "page": None}]

    if extension == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return [
            {"text": page.extract_text() or "", "page": page_number}
            for page_number, page in enumerate(reader.pages, start=1)
        ]

    if extension == ".docx":
        from docx import Document

        document = Document(str(path))
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                parts.append(" ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
        return [{"text": "\n".join(parts), "page": None}]

    if extension == ".doc":
        return [{"text": extract_legacy_office_text(path, extension), "page": None}]

    if extension == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        sections = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
            sections.append({"text": "\n".join(parts), "page": slide_number})
        return sections

    if extension == ".ppt":
        return [{"text": extract_legacy_office_text(path, extension), "page": None}]

    raise ValueError(f"Unsupported file type: {extension}")


def extract_legacy_office_text(path: Path, extension: str) -> str:
    if extension == ".doc" and shutil.which("textutil"):
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout

    if extension == ".doc" and shutil.which("antiword"):
        result = subprocess.run(
            ["antiword", str(path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout

    if extension == ".ppt" and shutil.which("catppt"):
        result = subprocess.run(
            ["catppt", str(path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout

    raise ValueError(
        f"Old {extension} files need a system converter. Please save the file as "
        f"{'.docx' if extension == '.doc' else '.pptx'} and upload again."
    )


def generate_answer(question: str):
    raise RuntimeError("External AI answers are disabled. Upload documents and answer from retrieved document context only.")


def decode_wav_pcm16(audio_bytes: bytes) -> tuple[bytes, int]:
    with wave.open(io.BytesIO(audio_bytes), "rb") as wav_file:
        channels = wav_file.getnchannels()
        sample_width = wav_file.getsampwidth()
        sample_rate = wav_file.getframerate()
        frames = wav_file.readframes(wav_file.getnframes())

    if sample_width != 2:
        raise ValueError("Streaming STT expects 16-bit PCM WAV chunks.")

    if channels == 1:
        return frames, sample_rate

    if channels != 2:
        raise ValueError("Streaming STT supports mono or stereo WAV chunks.")

    mono = bytearray(len(frames) // 2)
    write_offset = 0
    for offset in range(0, len(frames), 4):
        left = int.from_bytes(frames[offset : offset + 2], "little", signed=True)
        right = int.from_bytes(frames[offset + 2 : offset + 4], "little", signed=True)
        sample = int((left + right) / 2)
        mono[write_offset : write_offset + 2] = sample.to_bytes(2, "little", signed=True)
        write_offset += 2
    return bytes(mono), sample_rate


def has_gemini_stt_key() -> bool:
    return bool(os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY"))


def transcribe_audio_with_gemini(audio_bytes: bytes, mime_type: str = "audio/wav") -> dict:
    global GEMINI_STT_CLIENT
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is required for Gemini meeting transcription.")

    if GEMINI_STT_CLIENT is None:
        from google import genai

        GEMINI_STT_CLIENT = genai.Client(api_key=api_key)

    from google.genai import types

    model = os.getenv("GEMINI_STT_MODEL", DEFAULT_GEMINI_STT_MODEL)
    response = GEMINI_STT_CLIENT.models.generate_content(
        model=model,
        contents=[
            "Transcribe only the clearly spoken English words in this meeting audio. "
            "Return only the transcript text. If no speech is clear, return an empty string.",
            types.Part.from_bytes(data=audio_bytes, mime_type=mime_type or "audio/wav"),
        ],
    )
    transcript = clean_transcript((getattr(response, "text", "") or "").strip())
    return {
        "speechDetected": bool(transcript),
        "language": "english" if transcript else "unknown",
        "iso639_1": "en" if transcript else "unknown",
        "languageConfidence": 0.9 if transcript else 0.0,
        "transcript": transcript,
        "transcriptionConfidence": 0.85 if transcript else 0.0,
        "meaningful": is_meaningful_question_or_request(transcript) if transcript else False,
        "model": f"gemini/{model}",
        "isFinal": True,
    }


class VoskStreamingSession:
    def __init__(self, language_code: str = "en"):
        try:
            from vosk import KaldiRecognizer, Model
        except ImportError as exc:
            raise RuntimeError("Vosk is not installed for streaming speech-to-text.") from exc

        global STREAMING_VOSK_MODEL
        if ensure_vosk_model_path is None:
            raise RuntimeError("Vosk model path helper is unavailable.")

        model_path = ensure_vosk_model_path()
        if not model_path.exists():
            raise RuntimeError(f"Vosk model not found: {model_path}")

        if STREAMING_VOSK_MODEL is None:
            STREAMING_VOSK_MODEL = Model(str(model_path))

        self.model = STREAMING_VOSK_MODEL
        self.recognizer = KaldiRecognizer(self.model, 16000)
        self.recognizer.SetWords(False)
        self.language_code = language_code or "en"
        self.last_partial = ""
        self.last_used_at = time.time()

    def accept_chunk(self, audio_bytes: bytes) -> dict:
        pcm, sample_rate = decode_wav_pcm16(audio_bytes)
        if sample_rate != 16000:
            raise ValueError(f"Expected 16000 Hz WAV chunks, got {sample_rate} Hz.")

        self.last_used_at = time.time()
        if self.recognizer.AcceptWaveform(pcm):
            text = json.loads(self.recognizer.Result()).get("text", "").strip()
            self.last_partial = ""
            return {
                "speechDetected": bool(text),
                "language": "english",
                "iso639_1": "en",
                "languageConfidence": 0.95 if text else 0.0,
                "transcript": text,
                "transcriptionConfidence": 0.9 if text else 0.0,
                "meaningful": is_meaningful_question_or_request(clean_transcript(text)) if text else False,
                "model": "vosk/streaming",
                "isFinal": True,
            }

        text = json.loads(self.recognizer.PartialResult()).get("partial", "").strip()
        if text == self.last_partial:
            text = ""
        else:
            self.last_partial = text
        return {
            "speechDetected": bool(text),
            "language": "english" if text else "unknown",
            "iso639_1": "en" if text else "unknown",
            "languageConfidence": 0.7 if text else 0.0,
            "transcript": text,
            "transcriptionConfidence": 0.65 if text else 0.0,
            "meaningful": is_meaningful_question_or_request(clean_transcript(text)) if text else False,
            "model": "vosk/streaming-partial",
            "isFinal": False,
        }

    def finish(self) -> dict:
        self.last_used_at = time.time()
        text = json.loads(self.recognizer.FinalResult()).get("text", "").strip()
        self.last_partial = ""
        return {
            "speechDetected": bool(text),
            "language": "english" if text else "unknown",
            "iso639_1": "en" if text else "unknown",
            "languageConfidence": 0.95 if text else 0.0,
            "transcript": text,
            "transcriptionConfidence": 0.9 if text else 0.0,
            "meaningful": is_meaningful_question_or_request(clean_transcript(text)) if text else False,
            "model": "vosk/streaming-final",
            "isFinal": True,
        }


def prune_streaming_stt_sessions(now: Optional[float] = None) -> None:
    now = now or time.time()
    expired = [
        session_id
        for session_id, session in STREAMING_STT_SESSIONS.items()
        if now - session.last_used_at > STREAMING_STT_TTL_SECONDS
    ]
    for session_id in expired:
        STREAMING_STT_SESSIONS.pop(session_id, None)


def should_use_streaming_stt(requested_language: str, mime_type: str, session_id: str) -> bool:
    if not session_id:
        return False
    if os.getenv("ORBYNE_STREAMING_STT", "auto").strip().lower() in {"0", "false", "off", "disabled"}:
        return False
    if requested_language not in {"", "auto", "en"}:
        return False
    return "wav" in mime_type.lower()


def transcribe_streaming_audio(audio_bytes: bytes, session_id: str, final_chunk: bool = False):
    with STREAMING_STT_LOCK:
        prune_streaming_stt_sessions()
        session = STREAMING_STT_SESSIONS.get(session_id)
        if session is None:
            session = VoskStreamingSession()
            STREAMING_STT_SESSIONS[session_id] = session

        result = session.accept_chunk(audio_bytes)
        if final_chunk:
            final_result = session.finish()
            STREAMING_STT_SESSIONS.pop(session_id, None)
            if final_result["transcript"]:
                result = final_result
        return result


@app.get("/health")
def health():
    return {
        "status": "ok",
        "provider": "documents",
        "externalAiEnabled": False,
        "speechToText": "vosk/streaming",
    }


def record_audio_debug(entry: dict):
    AUDIO_DEBUG_LOGS.append(entry)
    logger.info("audio_chunk %s", json.dumps(entry, ensure_ascii=True, default=str))


def transcript_response(
    *,
    transcript: str,
    model: str,
    chunk_number: int,
    chunk_size: int,
    duration: float,
    sample_rate: int,
    audio_energy: float,
    voice_activity: bool,
    language: str = "unknown",
    language_confidence: float = 0.0,
    transcription_confidence: float = 0.0,
    meaningful: bool = False,
    discarded: bool = False,
    reason: str = "",
    is_final: bool = True,
    stage_timings: Optional[dict[str, float]] = None,
    started_at: float,
):
    latency_ms = int((time.time() - started_at) * 1000)
    merged_stage_timings = dict(stage_timings or {})
    merged_stage_timings.setdefault("totalBackendMs", latency_ms)
    response = TranscriptResponse(
        transcript=transcript,
        model=model,
        isFinal=is_final,
        chunkNumber=chunk_number,
        chunkSize=chunk_size,
        duration=round(duration, 3),
        sampleRate=sample_rate,
        audioEnergy=round(audio_energy, 6),
        voiceActivity=voice_activity,
        language=language,
        languageConfidence=round(language_confidence, 3),
        transcriptionConfidence=round(transcription_confidence, 3),
        meaningful=meaningful,
        discarded=discarded,
        reason=reason,
        latencyMs=latency_ms,
        stageTimings=merged_stage_timings,
    )
    record_audio_debug(response.dict())
    return response


@app.get("/debug/audio")
def audio_debug():
    return {"chunks": list(AUDIO_DEBUG_LOGS)}


@app.get("/")
def website():
    return FileResponse(BASE_DIR / "index.html")


@app.get("/icon.png")
def icon():
    return FileResponse(BASE_DIR / "icon.png")


@app.get("/logo.png")
def logo():
    return FileResponse(BASE_DIR / "logo.png")


@app.post("/answer", response_model=AnswerResponse)
def answer(request: AnswerRequest):
    try:
        answer_text, model = generate_answer(request.question.strip())
        return AnswerResponse(answer=answer_text, model=model)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@app.post("/answer-stream")
def answer_stream(request: AnswerRequest):
    raise HTTPException(
        status_code=410,
        detail="External AI answer streaming is disabled. Answers must come from uploaded documents only.",
    )


@app.post("/knowledge", response_model=KnowledgeResponse)
async def upload_knowledge(file: UploadFile = File(...)):
    filename = file.filename or "knowledge"
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_UPLOAD_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_UPLOAD_EXTENSIONS))
        raise HTTPException(status_code=400, detail=f"Unsupported file type. Supported: {supported}")

    suffix = extension or ".txt"
    temp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_path = Path(temp_file.name)
            while chunk := await file.read(1024 * 1024):
                temp_file.write(chunk)

        sections = [
            {"text": normalize_text(section["text"]), "page": section["page"]}
            for section in extract_upload_sections(temp_path, extension)
        ]
        sections = [section for section in sections if section["text"]]
        if not sections:
            raise HTTPException(status_code=400, detail="No readable text found in this file.")

        chunks = [
            KnowledgeChunk(text=chunk, filename=filename, page=section["page"])
            for section in sections
            for chunk in stream_chunks(section["text"])
        ]
        return KnowledgeResponse(chunks=chunks, chunkCount=len(chunks), filename=filename)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"File processing error: {exc}") from exc
    finally:
        if temp_path and temp_path.exists():
            temp_path.unlink(missing_ok=True)


@app.post("/transcribe-audio", response_model=TranscriptResponse)
async def transcribe_meeting_audio(
    file: UploadFile = File(...),
    language: str = Form("auto"),
    chunk_number: int = Form(0),
    session_id: str = Form(""),
    final_chunk: bool = Form(False),
    duration: float = Form(0.0),
    sample_rate: int = Form(0),
    audio_energy: float = Form(0.0),
    voice_activity: bool = Form(True),
):
    started_at = time.time()
    read_started_at = time.time()
    audio_bytes = await file.read()
    stage_timings = {
        "audioReadMs": round((time.time() - read_started_at) * 1000, 2),
        "languageDetectionMs": 0.0,
        "speechToTextMs": 0.0,
        "streamingSttMs": 0.0,
    }
    if not audio_bytes:
        raise HTTPException(status_code=400, detail="Empty audio chunk.")

    mime_type = file.content_type or "audio/webm"
    try:
        requested_language = (language or "auto").strip().lower()
        session_id = (session_id or "").strip()
        transcript = ""
        confidence = 0.0
        model = ""
        language_code = "unknown"
        language_confidence = 0.0
        speech_detected = False
        meaningful = False
        transcribed = False

        if not should_use_streaming_stt(requested_language, mime_type, session_id):
            raise RuntimeError("Meeting audio speech-to-text requires sessioned 16 kHz WAV chunks.")

        stt_started_at = time.time()
        meeting_stt_provider = os.getenv("ORBYNE_MEETING_STT_PROVIDER", "vosk").strip().lower()
        if meeting_stt_provider == "gemini" and has_gemini_stt_key():
            try:
                result = transcribe_audio_with_gemini(audio_bytes, "audio/wav")
            except Exception as gemini_exc:
                logger.warning("Gemini meeting transcription failed; falling back to Vosk: %s", gemini_exc)
                result = transcribe_streaming_audio(audio_bytes, session_id, final_chunk=final_chunk)
                result["model"] = f'{result["model"]}+gemini-fallback'
        elif meeting_stt_provider == "gemini":
            result = transcribe_streaming_audio(audio_bytes, session_id, final_chunk=final_chunk)
            result["model"] = f'{result["model"]}+gemini-unavailable'
        else:
            result = transcribe_streaming_audio(audio_bytes, session_id, final_chunk=final_chunk)

            if (
                meeting_stt_provider in {"auto", "vosk-gemini", "hybrid"}
                and voice_activity
                and not (result.get("transcript") or "").strip()
                and has_gemini_stt_key()
            ):
                try:
                    result = transcribe_audio_with_gemini(audio_bytes, "audio/wav")
                except Exception as gemini_exc:
                    logger.warning("Gemini hybrid transcription fallback failed: %s", gemini_exc)
                    result["model"] = f'{result["model"]}+gemini-fallback-unavailable'
        streaming_stt_ms = round((time.time() - stt_started_at) * 1000, 2)
        stage_timings["speechToTextMs"] = streaming_stt_ms
        stage_timings["streamingSttMs"] = streaming_stt_ms
        language_code = result["iso639_1"]
        language_confidence = result["languageConfidence"]
        model = result["model"]
        transcript = result["transcript"]
        confidence = result["transcriptionConfidence"]
        speech_detected = bool(result["speechDetected"])
        meaningful = bool(result["meaningful"])
        transcribed = True
        is_final = bool(result.get("isFinal", True))

        has_transcript_text = bool((transcript or "").strip())
        if not speech_detected and not has_transcript_text:
            return transcript_response(
                transcript="",
                model=model,
                chunk_number=chunk_number,
                chunk_size=len(audio_bytes),
                duration=duration,
                sample_rate=sample_rate,
                audio_energy=audio_energy,
                voice_activity=voice_activity,
                language=language_code,
                language_confidence=language_confidence,
                discarded=True,
                reason="no_clear_speech",
                stage_timings=stage_timings,
                started_at=started_at,
            )
        if len(transcript.split()) > 24:
            return transcript_response(
                transcript="",
                model=model,
                chunk_number=chunk_number,
                chunk_size=len(audio_bytes),
                duration=duration,
                sample_rate=sample_rate,
                audio_energy=audio_energy,
                voice_activity=voice_activity,
                language=language_code,
                language_confidence=language_confidence,
                transcription_confidence=confidence,
                discarded=True,
                reason="implausible_transcript",
                stage_timings=stage_timings,
                started_at=started_at,
            )
        if confidence < 0.45 and not has_transcript_text:
            return transcript_response(
                transcript="",
                model=model,
                chunk_number=chunk_number,
                chunk_size=len(audio_bytes),
                duration=duration,
                sample_rate=sample_rate,
                audio_energy=audio_energy,
                voice_activity=voice_activity,
                language=language_code,
                language_confidence=language_confidence,
                transcription_confidence=confidence,
                discarded=True,
                reason="low_transcription_confidence",
                stage_timings=stage_timings,
                started_at=started_at,
            )

        cleaned = clean_transcript(transcript)
        meaningful = meaningful or is_meaningful_question_or_request(cleaned)
        stage_timings["totalBackendMs"] = round((time.time() - started_at) * 1000, 2)

        return transcript_response(
            transcript=cleaned,
            model=model,
            chunk_number=chunk_number,
            chunk_size=len(audio_bytes),
            duration=duration,
            sample_rate=sample_rate,
            audio_energy=audio_energy,
            voice_activity=voice_activity,
            language=language_code,
            language_confidence=language_confidence,
            transcription_confidence=confidence,
            meaningful=meaningful,
            is_final=is_final,
            reason="",
            stage_timings=stage_timings,
            started_at=started_at,
        )
    except Exception as exc:
        logger.warning("Meeting transcription skipped after STT error: %s", exc)
        stage_timings["totalBackendMs"] = round((time.time() - started_at) * 1000, 2)
        return transcript_response(
            transcript="",
            model="speech-to-text",
            chunk_number=chunk_number,
            chunk_size=len(audio_bytes),
            duration=duration,
            sample_rate=sample_rate,
            audio_energy=audio_energy,
            voice_activity=voice_activity,
            discarded=True,
            reason="stt_unavailable",
            stage_timings=stage_timings,
            started_at=started_at,
        )


if __name__ == "__main__":
    import threading
    import webbrowser

    import uvicorn

    port = int(os.getenv("PORT", "8000"))
    url = f"http://127.0.0.1:{port}"
    print(f"ORBYNECUE web app: {url}")
    if os.getenv("ORBYNE_OPEN_BROWSER", "true").lower() != "false":
        threading.Timer(1.0, lambda: webbrowser.open(url)).start()
    uvicorn.run(app, host="127.0.0.1", port=port)
