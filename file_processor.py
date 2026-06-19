# file_processor.py - Production RAG Processor (PyInstaller-safe)

import os
import csv
import re
import numpy as np
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openai import OpenAI

from config import load_environment

# ======================================================
# Load .env correctly (works in EXE + normal Python)
# ======================================================
load_environment()
client = None

# ===============================
# GLOBAL DOCUMENT STORE
# ===============================
DOCUMENT_CHUNKS = []
DOCUMENT_EMBEDDINGS = []
DOCUMENT_METADATA = []

EMBEDDING_MODEL = "text-embedding-3-small"


def get_openai_client():
    global client
    if client is None:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise RuntimeError("OPENAI_API_KEY is required to process knowledge files.")
        client = OpenAI(api_key=api_key)
    return client

# ===============================
# STREAMING CHUNKER (BEST PRACTICE)
# ===============================
def stream_chunks(text, max_words=300):
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    buf = []
    word_count = 0

    for p in paragraphs:
        words = p.split()
        if not words:
            continue
        if word_count + len(words) <= max_words:
            buf.append(p)
            word_count += len(words)
        else:
            if buf:
                yield "\n".join(buf)
            buf = [p]
            word_count = len(words)

    if buf:
        yield "\n".join(buf)

# ===============================
# TEXT EXTRACTION
# ===============================
def extract_text(path):
    return "\n".join(item["text"] for item in extract_text_sections(path))


def extract_text_sections(path):
    ext = os.path.splitext(path)[1].lower()
    sections = []
    filename = os.path.basename(path)

    if ext == ".pdf":
        reader = PdfReader(path)
        for page_number, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text()
            if page_text:
                sections.append({"text": page_text, "filename": filename, "page": page_number})

    elif ext == ".docx":
        doc = Document(path)
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        sections.append({"text": text, "filename": filename, "page": None})

    elif ext == ".pptx":
        prs = Presentation(path)
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            sections.append({"text": slide_text, "filename": filename, "page": slide_number})

    elif ext == ".csv":
        with open(path, newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            text = ""
            for row in reader:
                text += " ".join(row) + "\n"
            sections.append({"text": text, "filename": filename, "page": None})

    else:
        raise ValueError("Unsupported file type")

    return sections

# ===============================
# NORMALIZATION
# ===============================
def normalize_text(text: str) -> str:
    text = re.sub(r"\n{2,}", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()

# ===============================
# MAIN ENTRY
# ===============================
def process_file(path):
    global DOCUMENT_CHUNKS, DOCUMENT_EMBEDDINGS, DOCUMENT_METADATA

    sections = extract_text_sections(path)

    DOCUMENT_CHUNKS.clear()
    DOCUMENT_EMBEDDINGS.clear()
    DOCUMENT_METADATA.clear()

    for section in sections:
        text = normalize_text(section["text"])
        for chunk in stream_chunks(text):
            DOCUMENT_CHUNKS.append(chunk)
            DOCUMENT_METADATA.append({"filename": section["filename"], "page": section["page"]})

    try:
        openai_client = get_openai_client()
        for chunk in DOCUMENT_CHUNKS:
            emb = openai_client.embeddings.create(
                model=EMBEDDING_MODEL,
                input=chunk
            ).data[0].embedding
            DOCUMENT_EMBEDDINGS.append(np.array(emb))
    except Exception as exc:
        DOCUMENT_EMBEDDINGS.clear()
        print(f"OpenAI embeddings unavailable; using local keyword search: {exc}")

    return len(DOCUMENT_CHUNKS)
